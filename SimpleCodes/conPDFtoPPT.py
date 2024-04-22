import os
from PyPDF2 import PdfReader
from pptx import Presentation

def pdf_to_text():
        pdf_file = "lmao.pdf"
        text = ""
        with open(pdf_file, 'rb') as f:
                reader = PdfReader(f)
                for page_num in range(len(reader.pages)):
                    page_text = reader.pages[page_num].extract_text()
                    text += page_text
        return text

def pdf_to_ppt(output_file):
        text = pdf_to_text()
	prs = Presentation()
	slides = text_split('\n\n')
	for slide_content in slides:
		slide = prs.slides.add_slides(prs.slide_layout[1])
		slide.shapes.title.text = slides_content
	prs.save(output_file)

#example usage 
output_ppt_file = "output_ppt.pptx"

pdf_to_ppt(output_ppt_file)






